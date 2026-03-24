#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generar_pptx.py — FifoCleanup v1.0 PowerPoint Generator
Replica PRESENTACION_TI.html con animaciones OOXML nativas
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
from PIL import Image, ImageDraw
import io, math, os

# ── Logos ─────────────────────────────────────────────────────────────
_BASE = r"c:\Users\IDC INGENIERIA\OneDrive\IDC\Proyectos\Ejecucion_Proyectos\03_Almacenamiento_FIFO"
LOGO_WHITE = os.path.join(_BASE, "idc-logo-white.png")  # para fondos oscuros
LOGO_DARK  = os.path.join(_BASE, "idc-logo-dark.png")   # para fondos claros

def add_logo(slide, dark_bg=True):
    """Agrega el logo IDC en la esquina inferior derecha."""
    logo_path = LOGO_WHITE if dark_bg else LOGO_DARK
    lw = Inches(1.3)
    lh = Inches(0.64)
    lx = SW - lw - Inches(0.3)
    ly = SH - lh - Inches(0.18)
    pic = slide.shapes.add_picture(logo_path, lx, ly, lw, lh)
    return pic

# ── Íconos Pillow (sin emoji) ─────────────────────────────────────────
_W  = (255, 255, 255, 255)   # blanco opaco
_T  = (255, 255, 255,  20)   # blanco semi-transparente (relleno interior)
_IC = 64                      # tamaño base de ícono en píxeles
_ICON_CACHE: dict = {}

def _mk(fn, size=_IC):
    img = Image.new('RGBA', (size, size), (0, 0, 0, 0))
    fn(ImageDraw.Draw(img), size)
    buf = io.BytesIO(); img.save(buf, 'PNG'); buf.seek(0)
    return buf

def _get_icon(name):
    if name not in _ICON_CACHE:
        if name not in _IFUNCS:
            raise ValueError(f"Icono desconocido: '{name}'")
        _ICON_CACHE[name] = _mk(_IFUNCS[name])
    _ICON_CACHE[name].seek(0)
    return _ICON_CACHE[name]

def add_icon(slide, name, x, y, w, h):
    """Inserta un ícono PNG (blanco sobre transparente) en el slide."""
    return slide.shapes.add_picture(_get_icon(name), x, y, w, h)

# ── Funciones de dibujo de íconos ─────────────────────────────────────
def _ico_warning(d, s):
    cx = s//2
    d.polygon([(cx, 6), (s-7, s-7), (7, s-7)], outline=_W, width=3)
    d.line([(cx, 18), (cx, s//2+4)], fill=_W, width=3)
    d.ellipse([cx-4, s//2+9, cx+4, s//2+17], fill=_W)

def _ico_bulb(d, s):
    cx = s//2; r = 15
    d.arc([cx-r, 4, cx+r, 4+r*2], start=0, end=360, fill=_W, width=3)
    bw = 8
    d.rectangle([cx-bw, 4+r*2-2, cx+bw, s-14], outline=_W, width=2)
    d.line([(cx-11, s-14), (cx+11, s-14)], fill=_W, width=2)
    d.line([(cx-7,  s-10), (cx+7,  s-10)], fill=_W, width=2)

def _ico_compare(d, s):
    m = 8; y1 = s//2-8; y2 = s//2+8
    d.line([(m, y1), (s-m, y1)], fill=_W, width=2)
    d.polygon([(s-m, y1), (s-m-8, y1-5), (s-m-8, y1+5)], fill=_W)
    d.line([(m, y2), (s-m, y2)], fill=_W, width=2)
    d.polygon([(m, y2), (m+8, y2-5), (m+8, y2+5)], fill=_W)

def _ico_layers(d, s):
    m = 8; h = 10; gap = 7
    for i in range(3):
        y = m + i*(h+gap)
        d.rectangle([m, y, s-m, y+h], outline=_W, width=2)

def _ico_cycle(d, s):
    m = 8
    d.arc([m, m, s-m, s-m], start=40, end=320, fill=_W, width=3)
    ax = int(s//2 + (s//2-m)*math.cos(math.radians(40)))
    ay = int(s//2 + (s//2-m)*math.sin(math.radians(40)))
    d.polygon([(ax, ay), (ax-8, ay-4), (ax-6, ay+8)], fill=_W)

def _ico_chart(d, s):
    m = 8
    d.line([(m, s-m), (s-m, s-m)], fill=_W, width=2)
    d.line([(m, m),   (m,   s-m)], fill=_W, width=2)
    for i, h in enumerate([0.45, 0.80, 0.60, 0.95]):
        x  = m + 6 + i*12
        bh = int((s-2*m-4)*h)
        d.rectangle([x, s-m-bh, x+9, s-m], fill=_W)

def _ico_bolt(d, s):
    pts = [(s//2+8, 5), (12, s//2+4), (s//2, s//2+4),
           (s//2-8, s-5), (s-12, s//2-4), (s//2, s//2-4)]
    d.polygon(pts, fill=_W)

def _ico_shield(d, s):
    cx = s//2
    d.polygon([(cx,6),(s-8,14),(s-8,s//2+6),(cx,s-6),(8,s//2+6),(8,14)],
              outline=_W, width=3)
    d.line([(cx-8, s//2), (cx-2, s//2+7), (cx+10, s//2-6)], fill=_W, width=3)

def _ico_check_circle(d, s):
    m = 6
    d.ellipse([m, m, s-m, s-m], outline=_W, width=3)
    cx, cy = s//2, s//2
    d.line([(cx-10, cy+1), (cx-3, cy+8), (cx+10, cy-7)], fill=_W, width=3)

def _ico_download(d, s):
    cx = s//2; m = 10
    d.line([(cx, m), (cx, s-m-8)], fill=_W, width=3)
    d.polygon([(cx, s-m-4), (cx-10, s-m-14), (cx+10, s-m-14)], fill=_W)
    d.line([(m, s-m), (s-m, s-m)], fill=_W, width=3)

def _ico_undo(d, s):
    m = 8
    d.arc([m+2, m, s-m-2, s//2+(s//2-m)], start=0, end=180, fill=_W, width=3)
    d.polygon([(m+2, s//2-2), (m+12, s//2-10), (m+12, s//2+2)], fill=_W)
    d.line([(s-m-2, s//2-2), (s-m-2, s-m)], fill=_W, width=3)
    d.line([(s//2, s-m), (s-m-2, s-m)], fill=_W, width=3)

def _ico_document(d, s):
    m = 8; fold = 14
    d.polygon([(m,m),(s-m-fold,m),(s-m,m+fold),(s-m,s-m),(m,s-m)],
              outline=_W, width=2)
    d.polygon([(s-m-fold,m),(s-m,m+fold),(s-m-fold,m+fold)], outline=_W, width=2)
    lm = m+6
    for i, y in enumerate([m+20, m+28, m+36, m+44]):
        end = s-lm-8 if i < 3 else s-lm-18
        d.line([(lm, y), (end, y)], fill=_W, width=2)

def _ico_speech(d, s):
    bh = int(s*0.55); m = 8
    cx = s//2
    pts = [(m,m),(s-m,m),(s-m,m+bh),(cx+6,m+bh),(cx,s-m),(cx-6,m+bh),(m,m+bh)]
    d.polygon(pts, outline=_W, width=2)
    cy = m + bh//2
    d.line([(m+8, cy-4), (s-m-8, cy-4)], fill=_W, width=2)
    d.line([(m+8, cy+4), (s-m-16, cy+4)], fill=_W, width=2)

def _ico_flag(d, s):
    px = 16
    d.line([(px, 8), (px, s-8)], fill=_W, width=3)
    d.polygon([(px,10),(s-10,22),(px,34)], fill=_W)

def _ico_gear(d, s):
    cx, cy = s//2, s//2
    ri = s//2-12; ro = s//2-6; rh = int(ri*0.45)
    d.ellipse([cx-ri,cy-ri,cx+ri,cy+ri], outline=_W, width=2)
    d.ellipse([cx-rh,cy-rh,cx+rh,cy+rh], outline=_W, width=2)
    for i in range(8):
        a = math.radians(i*45)
        x1=int(cx+ri*math.cos(a)); y1=int(cy+ri*math.sin(a))
        x2=int(cx+ro*math.cos(a)); y2=int(cy+ro*math.sin(a))
        d.line([(x1,y1),(x2,y2)], fill=_W, width=5)

def _ico_flask(d, s):
    cx = s//2; nw = 8; m = 10
    d.line([(cx-nw,m),(cx+nw,m)], fill=_W, width=2)
    d.line([(cx-nw,m),(cx-nw,m+12)], fill=_W, width=2)
    d.line([(cx+nw,m),(cx+nw,m+12)], fill=_W, width=2)
    bw = s//2-m
    d.polygon([(cx-nw,m+12),(cx-bw,s-m),(cx+bw,s-m),(cx+nw,m+12)],
              outline=_W, width=2)
    d.ellipse([cx-6,s-m-18,cx,  s-m-12], fill=_W)
    d.ellipse([cx+2,s-m-26,cx+8,s-m-20], fill=_W)

def _ico_play(d, s):
    m = 14
    d.polygon([(m,m),(s-m,s//2),(m,s-m)], fill=_W)

def _ico_log(d, s):
    m = 10; cw = 14
    d.rectangle([m, m+10, s-m, s-m], outline=_W, width=2)
    d.rectangle([s//2-cw//2,m-2,s//2+cw//2,m+14], outline=_W, width=2)
    for y in [m+22, m+31, m+40]:
        d.line([(m+6, y), (s-m-6, y)], fill=_W, width=2)

def _ico_no(d, s):
    m = 6; p = m+10
    d.ellipse([m,m,s-m,s-m], outline=_W, width=3)
    d.line([(p,p),(s-p,s-p)], fill=_W, width=3)

def _ico_thread(d, s):
    m = 8; bh = 10; gap = 8
    for i, f in enumerate([0.70, 0.45, 0.25]):
        y = m + i*(bh+gap)
        d.rectangle([m,y,s-m,y+bh], outline=_W, width=2)
        fw = int((s-2*m-4)*f)
        d.rectangle([m+2,y+2,m+2+fw,y+bh-2], fill=_W)

def _ico_pause(d, s):
    bw = 10; gap = 10; m = 16
    x0 = (s-2*bw-gap)//2
    d.rectangle([x0,m,x0+bw,s-m], fill=_W)
    d.rectangle([x0+bw+gap,m,x0+2*bw+gap,s-m], fill=_W)

def _ico_lock(d, s):
    cx = s//2; bw = 20; bh = 22; by = s//2+2
    d.rectangle([cx-bw//2,by,cx+bw//2,by+bh], outline=_W, width=2)
    d.arc([cx-12,s//4-2,cx+12,by+14], start=180, end=0, fill=_W, width=3)
    d.ellipse([cx-5,by+8,cx+5,by+18], outline=_W, width=2)

def _ico_close(d, s):
    m = 14
    d.line([(m,m),(s-m,s-m)], fill=_W, width=4)
    d.line([(s-m,m),(m,s-m)], fill=_W, width=4)

def _ico_trash(d, s):
    m = 10; cx = s//2
    d.line([(m,m+10),(s-m,m+10)], fill=_W, width=3)
    d.rectangle([cx-7,m+2,cx+7,m+10], outline=_W, width=2)
    bm = m+4
    d.polygon([(bm,m+10),(s-bm,m+10),(s-bm-4,s-m),(bm+4,s-m)], outline=_W, width=2)
    for dx in [-7, 0, 7]:
        d.line([(cx+dx,m+16),(cx+dx,s-m-5)], fill=_W, width=2)

_IFUNCS = {
    'warning':     _ico_warning,
    'bulb':        _ico_bulb,
    'compare':     _ico_compare,
    'layers':      _ico_layers,
    'cycle':       _ico_cycle,
    'chart':       _ico_chart,
    'bolt':        _ico_bolt,
    'shield':      _ico_shield,
    'check_circle':_ico_check_circle,
    'download':    _ico_download,
    'undo':        _ico_undo,
    'document':    _ico_document,
    'speech':      _ico_speech,
    'flag':        _ico_flag,
    'gear':        _ico_gear,
    'flask':       _ico_flask,
    'play':        _ico_play,
    'log':         _ico_log,
    'no':          _ico_no,
    'thread':      _ico_thread,
    'pause':       _ico_pause,
    'lock':        _ico_lock,
    'close':       _ico_close,
    'trash':       _ico_trash,
}

# ── Colores ──────────────────────────────────────────────────────────
B900 = RGBColor(0x0c, 0x19, 0x29)
B800 = RGBColor(0x0f, 0x24, 0x40)
B700 = RGBColor(0x16, 0x3a, 0x5f)
B600 = RGBColor(0x1a, 0x52, 0x76)
B500 = RGBColor(0x24, 0x71, 0xa3)
B400 = RGBColor(0x5d, 0xad, 0xe2)
B300 = RGBColor(0x85, 0xc1, 0xe9)
B200 = RGBColor(0xae, 0xd6, 0xf1)
B100 = RGBColor(0xd6, 0xea, 0xf8)
G500 = RGBColor(0x27, 0xae, 0x60)
G400 = RGBColor(0x2e, 0xcc, 0x71)
G300 = RGBColor(0x82, 0xe0, 0xaa)
G100 = RGBColor(0xd5, 0xf5, 0xe3)
R500 = RGBColor(0xe7, 0x4c, 0x3c)
R400 = RGBColor(0xec, 0x70, 0x63)
R100 = RGBColor(0xfa, 0xdb, 0xd8)
A500 = RGBColor(0xf3, 0x9c, 0x12)
A100 = RGBColor(0xfe, 0xf9, 0xe7)
GR700 = RGBColor(0x4a, 0x4a, 0x5a)
GR600 = RGBColor(0x6c, 0x6c, 0x80)
GR500 = RGBColor(0x8e, 0x8e, 0x9e)
GR400 = RGBColor(0xb0, 0xb0, 0xbe)
GR300 = RGBColor(0xd0, 0xd0, 0xda)
GR200 = RGBColor(0xe8, 0xe8, 0xec)
GR100 = RGBColor(0xf4, 0xf4, 0xf6)
WHITE = RGBColor(0xff, 0xff, 0xff)

# ── Dimensiones ──────────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)
ML = Inches(0.75)
MT = Inches(0.5)
CW = Inches(11.83)

# ── ID counter global ────────────────────────────────────────────────
_id = [100]

def nid():
    _id[0] += 1
    return _id[0]

# ── Helpers básicos ──────────────────────────────────────────────────

def new_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_bg_gradient(slide, color1, color2):
    """Simula gradiente con color1 (aproximación sólida)."""
    set_bg(slide, color1)

def txb(slide, text, x, y, w, h,
        size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False,
        name='Calibri', wrap=True, spacing_after=0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing_after:
        p.space_after = Pt(spacing_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = name
    r.font.color.rgb = color
    return box

def txb_multi(slide, lines, x, y, w, h,
              size=13, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, name='Calibri',
              line_spacing=1.0):
    """lines: list of str or dict {text, size, bold, color, name}"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(line, str):
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = name
            r.font.color.rgb = color
        elif isinstance(line, dict):
            r = p.add_run()
            r.text = line.get('text', '')
            r.font.size = Pt(line.get('size', size))
            r.font.bold = line.get('bold', bold)
            r.font.name = line.get('name', name)
            r.font.color.rgb = line.get('color', color)
    return box

def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def add_label(slide, text, x=None, y=None, color=B400):
    x = x or ML
    y = y or MT
    return txb(slide, text.upper(), x, y, Inches(5), Inches(0.3),
               size=8, bold=True, color=color, name='Calibri')

def add_title(slide, text, x=None, y=None, w=None, size=32, color=WHITE):
    x = x or ML
    y = y or (MT + Inches(0.32))
    w = w or CW
    return txb(slide, text, x, y, w, Inches(1.3),
               size=size, bold=True, color=color, name='Calibri Light')

def add_subtitle(slide, text, x=None, y=None, w=None, size=15, color=GR400):
    x = x or ML
    y = y or (MT + Inches(1.55))
    w = w or Inches(9)
    return txb(slide, text, x, y, w, Inches(0.9),
               size=size, bold=False, color=color, name='Calibri')

def card(slide, x, y, w, h, fill=WHITE, border=GR200, border_w=Pt(1)):
    return rect(slide, x, y, w, h, fill=fill, line_color=border, line_w=border_w)

def dark_card(slide, x, y, w, h):
    return rect(slide, x, y, w, h, fill=B800, line_color=B700, line_w=Pt(1))

# ── Transiciones ──────────────────────────────────────────────────────

def add_transition(slide):
    """Push hacia la izquierda (entre slides de contenido)."""
    trans = etree.SubElement(slide._element, qn('p:transition'))
    trans.set('spd', 'med')
    push = etree.SubElement(trans, qn('p:push'))
    push.set('dir', 'l')

def add_transition_fade(slide):
    """Fade — aparece con fundido suave (portada → índice)."""
    trans = etree.SubElement(slide._element, qn('p:transition'))
    trans.set('spd', 'med')
    etree.SubElement(trans, qn('p:fade'))

def add_transition_zoom(slide):
    """Zoom-in — crece desde el centro, como abrir una ventana (índice → contenido)."""
    trans = etree.SubElement(slide._element, qn('p:transition'))
    trans.set('spd', 'med')
    zoom = etree.SubElement(trans, qn('p:zoom'))
    zoom.set('dir', 'in')

def add_transition_zoom_out(slide):
    """Zoom-out: diapositiva aparece alejándose — regreso al menú/índice."""
    trans = etree.SubElement(slide._element, qn('p:transition'))
    trans.set('spd', 'med')
    zoom = etree.SubElement(trans, qn('p:zoom'))
    zoom.set('dir', 'out')

# ── Animaciones ───────────────────────────────────────────────────────

def add_animations(slide, shapes_delays):
    """
    shapes_delays: list of (shape, delay_ms)
    Agrega animación Fade-in auto-start.
    Las formas NO se ocultan inicialmente (sin p:set ni p:bldLst)
    para garantizar visibilidad aunque la animación no dispare.
    """
    if not shapes_delays:
        return

    slide_el = slide._element
    for old in slide_el.findall(qn('p:timing')):
        slide_el.remove(old)

    root_id = nid()
    auto_id = nid()

    shape_blocks = []
    for shape, delay in shapes_delays:
        spid = shape.shape_id
        tid1 = nid()
        tid2 = nid()
        grp = len(shape_blocks)
        # Sin p:set visibility -> la forma empieza VISIBLE siempre
        # p:animEffect fade juega como efecto visual encima
        shape_blocks.append(f"""
        <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:cTn id="{tid1}" presetID="10" presetClass="entr" presetSubtype="0"
                 fill="hold" grpId="{grp}" nodeType="withEffect">
            <p:stCondLst>
              <p:cond delay="{delay}"/>
            </p:stCondLst>
            <p:childTnLst>
              <p:animEffect transition="in" filter="fade">
                <p:cBhvr>
                  <p:cTn id="{tid2}" dur="600"/>
                  <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                </p:cBhvr>
              </p:animEffect>
            </p:childTnLst>
          </p:cTn>
        </p:par>""")

    shapes_xml = '\n'.join(shape_blocks)

    # Sin p:bldLst -> PowerPoint no oculta las formas antes de la animación
    timing_xml = f"""<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="{root_id}" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:par>
            <p:cTn id="{auto_id}" fill="hold">
              <p:stCondLst>
                <p:cond delay="0"/>
              </p:stCondLst>
              <p:childTnLst>
                {shapes_xml}
              </p:childTnLst>
            </p:cTn>
          </p:par>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst/>
</p:timing>"""

    timing_el = parse_xml(timing_xml)
    slide_el.append(timing_el)


# ═══════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════

_INDEX_TOPICS = [
    ('warning',      '01', 'El Problema'),
    ('bulb',         '02', 'La Solución'),
    ('compare',      '03', 'Antes vs Después'),
    ('layers',       '04', 'Arquitectura'),
    ('cycle',        '05', 'FIFO Dual'),
    ('chart',        '06', 'Dashboard'),
    ('bolt',         '07', 'Impacto'),
    ('shield',       '08', 'Seguridad'),
    ('check_circle', '09', 'Pruebas'),
    ('download',     '10', 'Instalación'),
    ('undo',         '11', 'Rollback'),
    ('document',     '12', 'Documentación'),
    ('speech',       '13', 'Solicitud'),
    ('flag',         '14', 'Resumen'),
]

def slide_index(prs, first=True, highlight=None):
    """
    Índice Visual — menú tipo OS con 14 tarjetas.
    first=True  → entrada con fade (portada → índice inicial, tarjetas animadas).
    first=False → retorno con zoom-out (contenido → índice), tarjetas visibles sin animación.
    highlight   → índice 0-13 de la tarjeta siguiente (resaltada en B400/blanco).
    """
    sl = new_slide(prs)
    set_bg(sl, B900)
    anims = []

    lb = add_label(sl, 'FifoCleanup v1.0', color=B400)
    tt = add_title(sl, 'Índice de la Presentación', size=30, color=WHITE)
    if first:
        anims.extend([(lb, 0), (tt, 80)])

    cols    = 5
    card_w  = Inches(2.23)
    card_h  = Inches(1.62)
    gap_x   = Inches(0.16)
    gap_y   = Inches(0.15)
    start_x = ML
    start_y = MT + Inches(1.3)

    for i, (icon, num, name) in enumerate(_INDEX_TOPICS):
        row = i // cols
        col = i % cols

        if row == 2:                             # última fila: 4 tarjetas centradas
            n_last = len(_INDEX_TOPICS) - cols * 2
            row_w  = n_last * card_w + (n_last - 1) * gap_x
            cx = (SW - row_w) / 2 + col * (card_w + gap_x)
        else:
            cx = start_x + col * (card_w + gap_x)

        cy = start_y + row * (card_h + gap_y)

        # Color según estado: normal / resaltado (próxima tarjeta)
        is_next   = (highlight is not None) and (i == highlight)
        card_fill = B700 if is_next else B800
        card_brd  = B400 if is_next else B600
        brd_w     = Pt(2) if is_next else Pt(1)
        num_fill  = B400 if is_next else B500
        txt_color = WHITE

        # Tarjeta base
        c = rect(sl, cx, cy, card_w, card_h,
                 fill=card_fill, line_color=card_brd, line_w=brd_w)

        # Badge número (esquina superior izquierda)
        nbg = rect(sl, cx + Inches(0.13), cy + Inches(0.12),
                   Inches(0.38), Inches(0.25), fill=num_fill)
        nt  = txb(sl, num, cx + Inches(0.13), cy + Inches(0.12),
                  Inches(0.38), Inches(0.24),
                  size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Ícono SVG centrado (PNG Pillow, sin emoji)
        ico_sz = Inches(0.45)
        ico_x  = cx + (card_w - ico_sz) / 2
        ico_y  = cy + Inches(0.38)
        it = add_icon(sl, icon, ico_x, ico_y, ico_sz, ico_sz)

        # Título inferior
        ttt = txb(sl, name, cx + Inches(0.1), cy + Inches(1.05),
                  card_w - Inches(0.2), Inches(0.45),
                  size=10, bold=True, color=txt_color, align=PP_ALIGN.CENTER)

        if first:
            delay = 250 + i * 70
            anims.extend([(c, delay), (nbg, delay), (nt, delay),
                          (it, delay), (ttt, delay)])

    add_logo(sl, dark_bg=True)

    if first:
        add_transition_fade(sl)      # portada → índice: fundido suave
        add_animations(sl, anims)
    else:
        add_transition_zoom_out(sl)  # contenido → índice: zoom-out (volver al menú)

    return sl

def slide_0_portada(prs):
    """Slide 0: Portada — bg-dark"""
    sl = new_slide(prs)
    set_bg(sl, B900)
    anims = []

    # Badge
    b = rect(sl, ML, MT, Inches(2.8), Inches(0.32), fill=B700)
    bt = txb(sl, 'SOLICITUD DE APROBACIÓN', ML + Inches(0.1), MT + Inches(0.04),
             Inches(2.6), Inches(0.25), size=8, bold=True, color=B400)
    anims.append((b, 0)); anims.append((bt, 0))

    # Title
    t = txb(sl, 'FifoCleanup  v1.0', ML, MT + Inches(0.45), Inches(9), Inches(1.3),
            size=52, bold=True, color=WHITE, name='Calibri Light')
    anims.append((t, 150))

    # Subtitle
    s = txb_multi(sl, [
        {'text': 'Gestión Automática de Almacenamiento FIFO', 'size': 20, 'color': B300, 'bold': False},
        {'text': 'para Servidor de Monitoreo Industrial', 'size': 20, 'color': GR400, 'bold': False},
    ], ML, MT + Inches(1.85), Inches(8), Inches(0.9), name='Calibri')
    anims.append((s, 300))

    # Info boxes
    info = [
        ('PRESENTADO POR', 'IDC Ingeniería'),
        ('PARA', 'Equipo de TI — ODL'),
        ('FECHA', 'Marzo 2026'),
        ('SERVIDOR DESTINO', 'SRVODLRTDNMICA'),
    ]
    x = ML
    y = MT + Inches(3.1)
    bw = Inches(2.6)
    for label_txt, val_txt in info:
        bg = rect(sl, x, y, bw, Inches(0.75), fill=B800, line_color=B700)
        lt = txb(sl, label_txt, x + Inches(0.15), y + Inches(0.05),
                 bw - Inches(0.3), Inches(0.25), size=7, color=GR500, bold=True)
        vt = txb(sl, val_txt, x + Inches(0.15), y + Inches(0.32),
                 bw - Inches(0.3), Inches(0.35), size=12, color=WHITE, bold=True)
        anims.extend([(bg, 500), (lt, 500), (vt, 500)])
        x += bw + Inches(0.2)

    add_logo(sl, dark_bg=True)
    add_transition(sl)
    add_animations(sl, anims)
    return sl


def slide_1_problema(prs):
    """Slide 1: El Problema — bg-white"""
    sl = new_slide(prs)
    set_bg(sl, WHITE)
    anims = []

    lb = add_label(sl, 'El Problema', color=B500)
    anims.append((lb, 0))

    tt = txb_multi(sl, [
        {'text': 'El servidor de monitoreo  ', 'size': 28, 'color': GR700, 'bold': True},
        {'text': 'se queda sin espacio', 'size': 28, 'color': R500, 'bold': True},
    ], ML, MT + Inches(0.32), CW, Inches(0.9), name='Calibri Light')
    anims.append((tt, 100))

    sub = txb(sl, 'El servidor SRVODLRTDNMICA recolecta datos de vibración y condición 24/7.\nCuando el disco se llena, el monitoreo se detiene.',
              ML, MT + Inches(1.3), Inches(9), Inches(0.7), size=13, color=GR600)
    anims.append((sub, 200))

    # --- Left: Impact cards ---
    ly = MT + Inches(2.2)
    lx = ML
    lw = Inches(5.5)

    impacts = [
        (R100, R500, 'Pérdida de monitoreo',
         'El software deja de recolectar datos. Punto ciego operacional.'),
        (A100, A500, 'Intervención manual',
         '30–60 min de análisis + borrado manual por un ingeniero.'),
        (A100, A500, 'Riesgo de error humano',
         'Eliminación manual puede borrar datos recientes o archivos del sistema.'),
    ]
    for i, (bg_c, border_c, title_t, desc_t) in enumerate(impacts):
        cy = ly + i * Inches(1.15)
        c = card(sl, lx, cy, lw, Inches(1.05), fill=bg_c, border=border_c)
        ct = txb(sl, title_t, lx + Inches(0.15), cy + Inches(0.08),
                 lw - Inches(0.3), Inches(0.35), size=12, bold=True, color=GR700)
        cd = txb(sl, desc_t, lx + Inches(0.15), cy + Inches(0.48),
                 lw - Inches(0.3), Inches(0.45), size=10, color=GR600)
        anims.extend([(c, 300 + i*150), (ct, 300 + i*150), (cd, 300 + i*150)])

    # --- Right: Filesystem tree ---
    rx = ML + Inches(6)
    rw = Inches(5.58)
    rc = card(sl, rx, MT + Inches(2.2), rw, Inches(2.2), fill=GR100, border=GR200)
    rt = txb_multi(sl, [
        {'text': 'D:\\MonitoringData\\', 'size': 10, 'color': B500, 'bold': True, 'name': 'Consolas'},
        {'text': '  ├── Asset001\\   ← bomba, motor…', 'size': 10, 'color': A500, 'name': 'Consolas'},
        {'text': '      ├── 00\\   ← variable de medición', 'size': 10, 'color': GR600, 'name': 'Consolas'},
        {'text': '          ├── E\\   ← eventos → crece diario', 'size': 10, 'color': GR600, 'name': 'Consolas'},
        {'text': '          └── F\\   ← frecuencia → crece diario', 'size': 10, 'color': GR600, 'name': 'Consolas'},
        {'text': '  ├── Asset002\\', 'size': 10, 'color': A500, 'name': 'Consolas'},
        {'text': '  └── … crecimiento continuo sin límite', 'size': 10, 'color': GR400, 'name': 'Consolas'},
    ], rx + Inches(0.2), MT + Inches(2.3), rw - Inches(0.4), Inches(2.0))
    anims.extend([(rc, 300), (rt, 350)])

    rexp = card(sl, rx, MT + Inches(4.55), rw, Inches(0.7), fill=WHITE, border=R400, border_w=Pt(2))
    rexpt = txb(sl, 'No existe mecanismo automático que gestione el ciclo de vida de estos datos.',
                rx + Inches(0.15), MT + Inches(4.63), rw - Inches(0.3), Inches(0.55),
                size=11, bold=True, color=GR700)
    anims.extend([(rexp, 500), (rexpt, 500)])

    add_logo(sl, dark_bg=False)
    add_transition_zoom(sl)    # índice → primer contenido: zoom-in (abrir ventana)
    add_animations(sl, anims)
    return sl


def slide_2_solucion(prs):
    """Slide 2: La Solución — bg-gradient (B700)"""
    sl = new_slide(prs)
    set_bg(sl, B700)
    anims = []

    lb = add_label(sl, 'La Solución', color=B300)
    anims.append((lb, 0))

    tt = add_title(sl, 'FifoCleanup: gestión de espacio automática e inteligente',
                   size=28, color=WHITE)
    anims.append((tt, 100))

    sub = add_subtitle(sl, 'Política FIFO: cuando el disco se acerca al umbral, se eliminan automáticamente\nlos datos más antiguos, preservando siempre los más recientes.',
                       size=14, color=B200)
    anims.append((sub, 200))

    features = [
        ('chart',   'Dashboard inteligente',
         'Semáforo visual, gráficas de uso, proyección de días hasta disco lleno.'),
        ('cycle',   'Limpieza automática dual',
         'RF-07 evalúa cada 6h. RF-08 monitorea en tiempo real. Dos capas.'),
        ('flask',   'Simulación sin riesgo',
         'Genera datos sintéticos para probar FIFO antes de producción.'),
        ('log',     'Bitácora de auditoría',
         'Registro inmutable: qué se borró, cuándo, por qué, cuánto se liberó.'),
        ('gear',    '17 parámetros configurables',
         'Umbrales, frecuencias, throttling, caps. Calibrado para SRVODLRTDNMICA.'),
        ('shield',  'Cero impacto al monitoreo',
         'Hilos de baja prioridad, throttling de I/O. Desinstalable en 3 min.'),
    ]

    cols = 3
    cw = Inches(3.7)
    ch = Inches(1.45)
    cx0 = ML
    cy0 = MT + Inches(2.45)
    gap = Inches(0.16)

    for i, (icon, title_t, desc_t) in enumerate(features):
        col = i % cols
        row = i // cols
        cx = cx0 + col * (cw + gap)
        cy = cy0 + row * (ch + gap)
        c = dark_card(sl, cx, cy, cw, ch)
        ic = add_icon(sl, icon, cx + Inches(0.12), cy + Inches(0.08), Inches(0.4), Inches(0.4))
        ti = txb(sl, title_t, cx + Inches(0.15), cy + Inches(0.5), cw - Inches(0.3),
                 Inches(0.35), size=12, bold=True, color=WHITE)
        de = txb(sl, desc_t, cx + Inches(0.15), cy + Inches(0.85), cw - Inches(0.3),
                 Inches(0.55), size=10, color=B300)
        d = 300 + i * 100
        anims.extend([(c, d), (ic, d), (ti, d), (de, d)])

    add_logo(sl, dark_bg=True)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_3_antes_despues(prs):
    """Slide 3: Antes vs. Después — bg-light"""
    sl = new_slide(prs)
    set_bg(sl, GR100)
    anims = []

    lb = add_label(sl, 'Resultado esperado', color=B500)
    anims.append((lb, 0))
    tt = add_title(sl, 'Antes vs. Después', size=30, color=GR700)
    anims.append((tt, 100))

    cw = Inches(5.6)
    ch_total = Inches(5.5)
    ly = MT + Inches(1.3)

    # LEFT — Sin FifoCleanup
    lx = ML
    lbg = card(sl, lx, ly, cw, ch_total, fill=WHITE, border=R400, border_w=Pt(2))
    anims.append((lbg, 200))
    lhead = rect(sl, lx, ly, cw, Inches(0.38), fill=R100, line_color=R400)
    lheadt = txb(sl, 'Sin FifoCleanup  (hoy)', lx + Inches(0.15), ly + Inches(0.06),
                 cw - Inches(0.3), Inches(0.28), size=11, bold=True, color=R500)
    anims.extend([(lhead, 200), (lheadt, 200)])

    left_items = [
        ('30–60', 'minutos por evento', 'Intervención manual: conectarse, analizar, eliminar.', R500),
        ('Cada 2–4 semanas', '', 'Frecuencia impredecible de disco lleno.', GR700),
        ('24–48 h-hombre/año', '', 'Tiempo del ingeniero dedicado solo a limpieza.', GR700),
        ('Visibilidad: ninguna', '', 'No se sabe el estado del disco hasta que falla.', GR700),
    ]
    for i, (big, small, desc, col) in enumerate(left_items):
        iy = ly + Inches(0.5) + i * Inches(1.2)
        ic = card(sl, lx + Inches(0.15), iy, cw - Inches(0.3), Inches(1.0),
                  fill=R100 if i == 0 else WHITE, border=GR200)
        ibt = txb(sl, big, lx + Inches(0.3), iy + Inches(0.08), cw - Inches(0.5),
                  Inches(0.45), size=22, bold=True, color=col)
        idt = txb(sl, desc, lx + Inches(0.3), iy + Inches(0.55), cw - Inches(0.5),
                  Inches(0.38), size=10, color=GR600)
        anims.extend([(ic, 300 + i*120), (ibt, 300 + i*120), (idt, 300 + i*120)])

    # RIGHT — Con FifoCleanup
    rx = ML + cw + Inches(0.63)
    rbg = card(sl, rx, ly, cw, ch_total, fill=WHITE, border=G500, border_w=Pt(2))
    anims.append((rbg, 200))
    rhead = rect(sl, rx, ly, cw, Inches(0.38), fill=G100, line_color=G500)
    rheadt = txb(sl, 'Con FifoCleanup', rx + Inches(0.15), ly + Inches(0.06),
                 cw - Inches(0.3), Inches(0.28), size=11, bold=True, color=G500)
    anims.extend([(rhead, 200), (rheadt, 200)])

    right_items = [
        ('0', 'minutos de intervención', 'Limpieza 100% automática. Cero participación humana.', G500),
        ('Cero incidentes', '', 'RF-07 + RF-08 previenen antes de que ocurra.', GR700),
        ('< 2 h-hombre/año', '', 'Solo verificación del Dashboard al inicio de turno.', GR700),
        ('Visibilidad: tiempo real', '', 'Semáforo + gráficas + proyección de días restantes.', GR700),
    ]
    for i, (big, small, desc, col) in enumerate(right_items):
        iy = ly + Inches(0.5) + i * Inches(1.2)
        ic = card(sl, rx + Inches(0.15), iy, cw - Inches(0.3), Inches(1.0),
                  fill=G100 if i == 0 else WHITE, border=G300 if i == 0 else GR200)
        ibt = txb(sl, big, rx + Inches(0.3), iy + Inches(0.08), cw - Inches(0.5),
                  Inches(0.45), size=22, bold=True, color=col)
        idt = txb(sl, desc, rx + Inches(0.3), iy + Inches(0.55), cw - Inches(0.5),
                  Inches(0.38), size=10, color=GR600)
        anims.extend([(ic, 300 + i*120), (ibt, 300 + i*120), (idt, 300 + i*120)])

    add_logo(sl, dark_bg=False)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_4_arquitectura(prs):
    """Slide 4: Stack tecnológico — bg-dark"""
    sl = new_slide(prs)
    set_bg(sl, B900)
    anims = []

    lb = add_label(sl, 'Arquitectura', color=B400)
    anims.append((lb, 0))
    tt = add_title(sl, 'Stack tecnológico', color=WHITE)
    anims.append((tt, 100))
    sub = add_subtitle(sl, '100% C# / .NET 8.0 — sin componentes nativos ni dependencias exóticas',
                       size=13, color=B300)
    anims.append((sub, 200))

    # Arch boxes (left column)
    ax = ML
    ay = MT + Inches(1.85)
    aw = Inches(6.5)

    arch_items = [
        (B600, 'FifoCleanup.UI', 'WPF Application · MVVM · CommunityToolkit.Mvvm',
         'Dashboard  ·  Configuración  ·  Simulación  ·  Ejecución  ·  Bitácora'),
        (B700, 'FifoCleanup.Engine', 'Class Library · System.Text.Json',
         'RF-01 Inventario  ·  RF-02 Config  ·  RF-03 Simulación  ·  RF-04 Limpieza  ·  RF-05 Bitácora  ·  RF-07  ·  RF-08'),
        (B800, 'Sistema de Archivos', '',
         'D:\\MonitoringData\\Asset*\\Variable*\\{E,F}\\YYYY\\MM\\DD\\'),
    ]
    for i, (fill_c, arch_t, sub_t, detail_t) in enumerate(arch_items):
        by = ay + i * (Inches(1.35) + Inches(0.18))
        bh = Inches(1.35)
        b = rect(sl, ax, by, aw, bh, fill=fill_c, line_color=B500, line_w=Pt(1))
        bt = txb(sl, arch_t, ax + Inches(0.2), by + Inches(0.1), aw - Inches(0.4),
                 Inches(0.35), size=14, bold=True, color=WHITE)
        if sub_t:
            bs = txb(sl, sub_t, ax + Inches(0.2), by + Inches(0.45), aw - Inches(0.4),
                     Inches(0.25), size=10, color=B300)
            anims.append((bs, 300 + i*150))
        bd = txb(sl, detail_t, ax + Inches(0.2), by + Inches(0.72), aw - Inches(0.4),
                 Inches(0.55), size=10, color=B200, name='Consolas' if i == 2 else 'Calibri')
        anims.extend([(b, 300 + i*150), (bt, 300 + i*150), (bd, 300 + i*150)])
        # Arrow
        if i < 2:
            arr = txb(sl, '▼', ax + aw / 2 - Inches(0.2),
                      by + bh + Inches(0.02), Inches(0.4), Inches(0.2),
                      size=12, color=B400, align=PP_ALIGN.CENTER)
            anims.append((arr, 300 + i*150))

    # Decisions (right column)
    dx = ML + Inches(6.8)
    dw = Inches(4.78)
    dh_label = txb(sl, 'DECISIONES CLAVE', dx, ay - Inches(0.35), dw, Inches(0.3),
                   size=8, bold=True, color=GR500)
    anims.append((dh_label, 200))

    decisions = [
        ('🎯', '100% C# / .NET 8.0', 'Stack unificado. Sin C++ nativo. Simplifica instalación y testing.'),
        ('🧩', 'Engine separado de UI', 'Motor FIFO testeable independiente. Reusable como servicio en v2.0.'),
        ('📄', 'JSON unificado', 'Un solo archivo fifo_config.json con 17 parámetros tipados.'),
        ('📋', 'Bitácora CSV inmutable', 'Append-only. Abierto con Excel. Rotación a 100 MB.'),
        ('🧵', 'Hilos de baja prioridad', 'Limpieza cede CPU al monitoreo. Throttling 50ms entre operaciones.'),
    ]
    for i, (icon, dt, dd) in enumerate(decisions):
        dy = ay + i * Inches(0.88)
        dc = dark_card(sl, dx, dy, dw, Inches(0.82))
        dit = txb(sl, f'{icon}  {dt}', dx + Inches(0.15), dy + Inches(0.08),
                  dw - Inches(0.3), Inches(0.3), size=11, bold=True, color=WHITE)
        did = txb(sl, dd, dx + Inches(0.15), dy + Inches(0.42), dw - Inches(0.3),
                  Inches(0.35), size=9, color=B300)
        anims.extend([(dc, 400 + i*120), (dit, 400 + i*120), (did, 400 + i*120)])

    add_logo(sl, dark_bg=True)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_5_fifo_dual(prs):
    """Slide 5: Estrategia FIFO Dual — bg-white"""
    sl = new_slide(prs)
    set_bg(sl, WHITE)
    anims = []

    lb = add_label(sl, 'Protección automática', color=B500)
    anims.append((lb, 0))
    tt = add_title(sl, 'Estrategia FIFO Dual: defensa en profundidad', size=28, color=GR700)
    anims.append((tt, 100))
    sub = add_subtitle(sl, 'Dos modos complementarios garantizan que ningún escenario pase desapercibido.',
                       size=13, color=GR600)
    anims.append((sub, 200))

    cw = Inches(5.6)
    ch = Inches(4.5)
    ly = MT + Inches(2.15)

    # RF-07
    lx = ML
    lc = card(sl, lx, ly, cw, ch, fill=B100, border=B500, border_w=Pt(3))
    anims.append((lc, 300))
    badge = rect(sl, lx + Inches(0.2), ly + Inches(0.2), Inches(0.65), Inches(0.28), fill=B500)
    badget = txb(sl, 'RF-07', lx + Inches(0.22), ly + Inches(0.22), Inches(0.6), Inches(0.25),
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    lt = txb(sl, 'Limpieza Programada', lx + Inches(0.2), ly + Inches(0.6), cw - Inches(0.4),
             Inches(0.5), size=18, bold=True, color=B600)
    ls = txb(sl, 'Evaluación global periódica', lx + Inches(0.2), ly + Inches(1.1),
             cw - Inches(0.4), Inches(0.3), size=11, color=GR600)
    anims.extend([(badge, 300), (badget, 300), (lt, 350), (ls, 350)])

    items_07 = [
        'Se ejecuta cada 6 horas',
        'Evalúa TODO el disco (todos los Assets)',
        'Proyecta con historial de 7 días',
        'Si proyección > umbral → limpieza general FIFO',
        'Elimina datos más antiguos de todos los Assets',
    ]
    for i, item in enumerate(items_07):
        it = txb(sl, f'—  {item}', lx + Inches(0.2), ly + Inches(1.5) + i * Inches(0.48),
                 cw - Inches(0.4), Inches(0.4), size=11, color=GR700)
        anims.append((it, 400 + i*80))
    ideal = card(sl, lx + Inches(0.2), ly + ch - Inches(0.75), cw - Inches(0.4),
                 Inches(0.55), fill=B100, border=B200)
    idealt = txb(sl, 'Ideal para: Crecimiento gradual y uniforme',
                 lx + Inches(0.35), ly + ch - Inches(0.67), cw - Inches(0.6), Inches(0.4),
                 size=10, color=B600, bold=True)
    anims.extend([(ideal, 500), (idealt, 500)])

    # RF-08
    rx = ML + cw + Inches(0.63)
    rc = card(sl, rx, ly, cw, ch, fill=G100, border=G500, border_w=Pt(3))
    anims.append((rc, 300))
    rbadge = rect(sl, rx + Inches(0.2), ly + Inches(0.2), Inches(0.65), Inches(0.28), fill=G500)
    rbadget = txb(sl, 'RF-08', rx + Inches(0.22), ly + Inches(0.22), Inches(0.6), Inches(0.25),
                  size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rt = txb(sl, 'Monitoreo Preventivo', rx + Inches(0.2), ly + Inches(0.6), cw - Inches(0.4),
             Inches(0.5), size=18, bold=True, color=G500)
    rs = txb(sl, 'Reacción local en tiempo real', rx + Inches(0.2), ly + Inches(1.1),
             cw - Inches(0.4), Inches(0.3), size=11, color=GR600)
    anims.extend([(rbadge, 300), (rbadget, 300), (rt, 350), (rs, 350)])

    items_08 = [
        'FileSystemWatcher detecta archivos nuevos',
        'Agrupa eventos cada 10 segundos',
        'Evalúa solo el Asset afectado',
        'Si se llenará en < 3 días → limpieza local',
        'Elimina N días antiguos del Asset específico',
    ]
    for i, item in enumerate(items_08):
        it = txb(sl, f'—  {item}', rx + Inches(0.2), ly + Inches(1.5) + i * Inches(0.48),
                 cw - Inches(0.4), Inches(0.4), size=11, color=GR700)
        anims.append((it, 400 + i*80))
    rideal = card(sl, rx + Inches(0.2), ly + ch - Inches(0.75), cw - Inches(0.4),
                  Inches(0.55), fill=G100, border=G300)
    ridealt = txb(sl, 'Ideal para: Picos repentinos de un equipo',
                  rx + Inches(0.35), ly + ch - Inches(0.67), cw - Inches(0.6), Inches(0.4),
                  size=10, color=G500, bold=True)
    anims.extend([(rideal, 500), (ridealt, 500)])

    add_logo(sl, dark_bg=False)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_6_dashboard(prs):
    """Slide 6: Dashboard — bg-light"""
    sl = new_slide(prs)
    set_bg(sl, GR100)
    anims = []

    lb = add_label(sl, 'Interfaz de usuario', color=B500)
    anims.append((lb, 0))
    tt = add_title(sl, 'Dashboard: estado del servidor de un vistazo', size=26, color=GR700)
    anims.append((tt, 100))

    # Dashboard mockup card
    dy = MT + Inches(1.3)
    dw = Inches(11.83)
    dh = Inches(3.6)
    dc = card(sl, ML, dy, dw, dh, fill=WHITE, border=GR200, border_w=Pt(2))
    anims.append((dc, 200))

    # Semaphore circle (green)
    sem = rect(sl, ML + Inches(0.3), dy + Inches(0.5), Inches(1.0), Inches(1.0), fill=G500)
    pct = txb(sl, '72%', ML + Inches(0.35), dy + Inches(1.6), Inches(0.9), Inches(0.35),
              size=16, bold=True, color=G500, align=PP_ALIGN.CENTER)
    uso = txb(sl, 'Uso actual', ML + Inches(0.3), dy + Inches(1.98), Inches(1.0), Inches(0.25),
              size=9, color=GR500, align=PP_ALIGN.CENTER)
    anims.extend([(sem, 300), (pct, 300), (uso, 300)])

    # Stats
    stats = [
        ('Assets detectados', '12'),
        ('Datos monitoreados', '1.4 TB'),
        ('Crecimiento', '2.1 GB/día'),
        ('Disco lleno en', '324 días'),
    ]
    sx = ML + Inches(1.6)
    for i, (label_t, val_t) in enumerate(stats):
        slx = sx + i * Inches(2.55)
        slt = txb(sl, label_t.upper(), slx, dy + Inches(0.4), Inches(2.4), Inches(0.25),
                  size=7, bold=True, color=GR500)
        slv = txb(sl, val_t, slx, dy + Inches(0.7), Inches(2.4), Inches(0.45),
                  size=22, bold=True, color=GR700 if i < 3 else G500)
        anims.extend([(slt, 350 + i*80), (slv, 350 + i*80)])

    # Progress bar
    pb_bg = rect(sl, ML + Inches(1.6), dy + Inches(1.4), Inches(9.8), Inches(0.2), fill=GR200)
    pb_fill = rect(sl, ML + Inches(1.6), dy + Inches(1.4), Inches(7.06), Inches(0.2), fill=G500)
    pb_l = txb(sl, 'Capacidad del disco', ML + Inches(1.6), dy + Inches(1.2), Inches(4), Inches(0.22),
               size=9, color=GR500)
    pb_r = txb(sl, '1.4 TB / 2.0 TB', ML + Inches(7), dy + Inches(1.2), Inches(2), Inches(0.22),
               size=9, bold=True, color=GR700, align=PP_ALIGN.RIGHT)
    anims.extend([(pb_bg, 400), (pb_fill, 450), (pb_l, 400), (pb_r, 400)])

    # Mini asset list
    sep = rect(sl, ML + Inches(0.2), dy + Inches(1.85), dw - Inches(0.4), Pt(1), fill=GR200)
    anims.append((sep, 500))
    assets = [('Asset001', 24, '340 GB'), ('Asset002', 19, '260 GB'), ('Asset003', 13, '180 GB')]
    for i, (name, pct_v, size_t) in enumerate(assets):
        ay2 = dy + Inches(2.05) + i * Inches(0.42)
        aname = txb(sl, name, ML + Inches(0.3), ay2, Inches(1.1), Inches(0.35),
                    size=10, color=GR600, name='Consolas')
        abg = rect(sl, ML + Inches(1.5), ay2 + Inches(0.08), Inches(8.5), Inches(0.2), fill=GR200)
        afill = rect(sl, ML + Inches(1.5), ay2 + Inches(0.08), Inches(8.5 * pct_v / 100), Inches(0.2), fill=B400)
        asize = txb(sl, size_t, ML + Inches(10.1), ay2, Inches(1.4), Inches(0.35),
                    size=10, bold=True, color=GR700, align=PP_ALIGN.RIGHT)
        anims.extend([(aname, 550 + i*80), (abg, 550 + i*80), (afill, 600 + i*80), (asize, 550 + i*80)])

    # 4 module cards at bottom
    modules = [('gear',  'Configuración', '17 parámetros'),
               ('flask', 'Simulación',    'Pruebas sin riesgo'),
               ('play',  'Ejecución',     'Manual + RF-07 + RF-08'),
               ('log',   'Bitácora',      'Auditoría + CSV')]
    mw = Inches(2.75)
    mx = ML
    my = dy + dh + Inches(0.2)
    for i, (icon, mtitle, mdesc) in enumerate(modules):
        mx_i = mx + i*(mw + Inches(0.12))
        mc = card(sl, mx_i, my, mw, Inches(1.0), fill=WHITE, border=GR200)
        mic = add_icon(sl, icon, mx_i + (mw - Inches(0.38))/2, my + Inches(0.08),
                       Inches(0.38), Inches(0.38))
        mt = txb(sl, mtitle, mx_i + Inches(0.15), my + Inches(0.52),
                 mw - Inches(0.3), Inches(0.28), size=11, bold=True, color=GR700, align=PP_ALIGN.CENTER)
        md = txb(sl, mdesc,  mx_i + Inches(0.15), my + Inches(0.73),
                 mw - Inches(0.3), Inches(0.25), size=9, color=GR500, align=PP_ALIGN.CENTER)
        anims.extend([(mc, 600+i*100),(mic, 600+i*100),(mt, 600+i*100),(md, 600+i*100)])

    add_logo(sl, dark_bg=False)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_7_impacto(prs):
    """Slide 7: Impacto mínimo — bg-white"""
    sl = new_slide(prs)
    set_bg(sl, WHITE)
    anims = []

    lb = add_label(sl, 'Análisis de impacto', color=B500)
    anims.append((lb, 0))
    tt = add_title(sl, 'Impacto mínimo en el servidor', size=28, color=GR700)
    anims.append((tt, 100))
    sub = add_subtitle(sl, 'Diseñado específicamente para no interferir con el software de monitoreo.',
                       size=13, color=GR600)
    anims.append((sub, 200))

    metrics = [
        ('< 5%', 'CPU promedio', 'De 36 threads disponibles', 5),
        ('< 200', 'MB de RAM', 'De ~8 GB disponibles', 13),
        ('< 50', 'MB en disco', 'Tamaño de instalación', 1),
    ]
    mw = Inches(3.7)
    mx = ML
    my = MT + Inches(1.85)
    for i, (big, label_t, note_t, bar_pct) in enumerate(metrics):
        mc = card(sl, mx + i*(mw + Inches(0.16)), my, mw, Inches(1.8), fill=WHITE, border=GR200)
        mbt = txb(sl, big, mx + i*(mw + Inches(0.16)) + Inches(0.2), my + Inches(0.15),
                  mw - Inches(0.4), Inches(0.6), size=32, bold=True, color=B500, align=PP_ALIGN.CENTER)
        mlt = txb(sl, label_t, mx + i*(mw + Inches(0.16)) + Inches(0.2), my + Inches(0.75),
                  mw - Inches(0.4), Inches(0.3), size=12, bold=True, color=GR700, align=PP_ALIGN.CENTER)
        barbg = rect(sl, mx + i*(mw + Inches(0.16)) + Inches(0.2), my + Inches(1.18),
                     mw - Inches(0.4), Inches(0.18), fill=GR200)
        barfill = rect(sl, mx + i*(mw + Inches(0.16)) + Inches(0.2), my + Inches(1.18),
                       (mw - Inches(0.4)) * bar_pct / 100, Inches(0.18), fill=B400)
        mnote = txb(sl, note_t, mx + i*(mw + Inches(0.16)) + Inches(0.2), my + Inches(1.45),
                    mw - Inches(0.4), Inches(0.25), size=9, color=GR500, align=PP_ALIGN.CENTER)
        d = 300 + i * 120
        anims.extend([(mc, d), (mbt, d), (mlt, d), (barbg, d+50), (barfill, d+100), (mnote, d)])

    # Protection mechanisms
    ph = txb(sl, 'Mecanismos de protección', ML, MT + Inches(3.9), CW, Inches(0.35),
             size=14, bold=True, color=GR700)
    anims.append((ph, 400))

    protections = [
        ('thread', 'Hilos BelowNormal', 'Limpieza cede CPU al monitoreo automáticamente'),
        ('pause',  'Throttling 50ms',   'Pausa entre cada eliminación para reducir picos de I/O'),
        ('lock',   'Cap 20%',           'Máximo 20% de datos eliminados por ejecución'),
        ('shield', '3 niveles de captura', 'La app nunca crashea. Registra el error y continúa.'),
    ]
    pw = Inches(5.75)
    for i, (icon, ptitle, pdesc) in enumerate(protections):
        col = i % 2
        row = i // 2
        px = ML + col * (pw + Inches(0.33))
        py = MT + Inches(4.35) + row * Inches(1.0)
        pc  = card(sl, px, py, pw, Inches(0.85), fill=GR100, border=GR200)
        pic = add_icon(sl, icon, px + Inches(0.12), py + Inches(0.08), Inches(0.36), Inches(0.36))
        pit = txb(sl, ptitle, px + Inches(0.55), py + Inches(0.1),
                  pw - Inches(0.7), Inches(0.3), size=12, bold=True, color=GR700)
        pid = txb(sl, pdesc,  px + Inches(0.15), py + Inches(0.45),
                  pw - Inches(0.3), Inches(0.33), size=10, color=GR600)
        d = 500 + i * 100
        anims.extend([(pc, d), (pic, d), (pit, d), (pid, d)])

    add_logo(sl, dark_bg=False)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_8_seguridad(prs):
    """Slide 8: Seguridad — bg-dark"""
    sl = new_slide(prs)
    set_bg(sl, B900)
    anims = []

    lb = add_label(sl, 'Seguridad', color=B400)
    anims.append((lb, 0))
    tt = add_title(sl, 'Análisis de seguridad para TI', color=WHITE)
    anims.append((tt, 100))
    sub = add_subtitle(sl, 'FifoCleanup no abre puertos, no accede a internet y no modifica el sistema operativo.',
                       size=13, color=B300)
    anims.append((sub, 200))

    # Left: Lo que NO hace
    lx = ML
    lw = Inches(5.5)
    lh_label = txb(sl, 'LO QUE NO HACE', lx, MT + Inches(2.0), lw, Inches(0.28),
                   size=8, bold=True, color=GR500)
    anims.append((lh_label, 300))

    no_items = [
        'No requiere conexión a internet',
        'No abre puertos de red',
        'No modifica el registro de Windows',
        'No instala drivers ni servicios del SO',
        'No requiere permisos de administrador',
        'No tiene telemetría ni actualizaciones automáticas',
    ]
    for i, item in enumerate(no_items):
        ic = dark_card(sl, lx, MT + Inches(2.35) + i * Inches(0.73), lw, Inches(0.62))
        it = txb(sl, f'—  {item}', lx + Inches(0.15), MT + Inches(2.43) + i * Inches(0.73),
                 lw - Inches(0.3), Inches(0.4), size=11, color=WHITE)
        anims.extend([(ic, 350 + i*80), (it, 350 + i*80)])

    # Right: Dependencias
    rx = ML + lw + Inches(0.83)
    rw = Inches(5.0)
    rh_label = txb(sl, 'DEPENDENCIAS', rx, MT + Inches(2.0), rw, Inches(0.28),
                   size=8, bold=True, color=GR500)
    anims.append((rh_label, 300))

    deps = [
        ('.NET 8.0 Runtime', 'Microsoft (LTS)', G500),
        ('CommunityToolkit.Mvvm', 'Microsoft OSS', G500),
        ('LiveChartsCore', 'NuGet Gallery', G500),
        ('System.Text.Json', 'Microsoft (.NET)', G500),
    ]
    for i, (comp, origin, risk_c) in enumerate(deps):
        dy = MT + Inches(2.35) + i * Inches(0.55)
        drow = dark_card(sl, rx, dy, rw, Inches(0.48))
        dct = txb(sl, comp, rx + Inches(0.15), dy + Inches(0.08), Inches(2.4), Inches(0.3),
                  size=10, color=WHITE, name='Consolas')
        dot = txb(sl, origin, rx + Inches(2.65), dy + Inches(0.1), Inches(1.5), Inches(0.3),
                  size=10, color=GR400)
        drisk = rect(sl, rx + rw - Inches(1.0), dy + Inches(0.1), Inches(0.8), Inches(0.28), fill=G500)
        driskt = txb(sl, 'Bajo', rx + rw - Inches(0.98), dy + Inches(0.12), Inches(0.75), Inches(0.22),
                     size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        anims.extend([(drow, 400 + i*80), (dct, 400 + i*80), (dot, 400 + i*80),
                      (drisk, 400 + i*80), (driskt, 400 + i*80)])

    # File access info
    fac = dark_card(sl, rx, MT + Inches(4.65), rw, Inches(1.6))
    fac.line.color.rgb = G400
    fat = txb(sl, '¿Qué archivos modifica?', rx + Inches(0.15), MT + Inches(4.78),
              rw - Inches(0.3), Inches(0.3), size=12, bold=True, color=G400)
    fad = txb_multi(sl, [
        'Lee/escribe:  fifo_config.json (configuración)',
        'Agrega:  bitacora/*.csv (auditoría)',
        'Elimina:  Solo carpetas YYYY/MM/DD dentro de Asset/Variable/{E,F}/',
    ], rx + Inches(0.15), MT + Inches(5.15), rw - Inches(0.3), Inches(1.0),
        size=10, color=B300, name='Calibri')
    anims.extend([(fac, 500), (fat, 500), (fad, 520)])

    add_logo(sl, dark_bg=True)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_9_pruebas(prs):
    """Slide 9: Tests — bg-light"""
    sl = new_slide(prs)
    set_bg(sl, GR100)
    anims = []

    lb = add_label(sl, 'Calidad', color=B500)
    anims.append((lb, 0))
    tt = add_title(sl, '89 tests ejecutados — 100% de aprobación', size=26, color=GR700)
    anims.append((tt, 100))

    # Big stat
    bsx = ML
    bsy = MT + Inches(1.45)
    big = txb(sl, '82', bsx, bsy, Inches(2.5), Inches(1.8),
              size=80, bold=True, color=G500, align=PP_ALIGN.CENTER)
    bsub = txb(sl, 'de 82 ejecutables', bsx, bsy + Inches(1.8), Inches(2.5), Inches(0.4),
               size=13, bold=True, color=GR700, align=PP_ALIGN.CENTER)
    bsub2 = txb(sl, '0 fallidos  ·  7 N/A (v2.0)', bsx, bsy + Inches(2.25), Inches(2.5), Inches(0.3),
                size=10, color=GR500, align=PP_ALIGN.CENTER)
    bbar_bg = rect(sl, bsx + Inches(0.3), bsy + Inches(2.7), Inches(1.9), Inches(0.2), fill=GR200)
    bbar_fill = rect(sl, bsx + Inches(0.3), bsy + Inches(2.7), Inches(1.9), Inches(0.2), fill=G500)
    anims.extend([(big, 200), (bsub, 250), (bsub2, 280), (bbar_bg, 300), (bbar_fill, 350)])

    # Test table
    tx = ML + Inches(3.0)
    tw = Inches(8.58)
    ty = MT + Inches(1.45)

    # Header
    thead = rect(sl, tx, ty, tw, Inches(0.4), fill=B800)
    cols_h = ['Req.', 'Descripción', 'Tests', 'Estado']
    col_ws = [Inches(1.0), Inches(5.5), Inches(1.0), Inches(1.0)]
    cx = tx
    for h, cw2 in zip(cols_h, col_ws):
        ht = txb(sl, h, cx + Inches(0.1), ty + Inches(0.08), cw2 - Inches(0.1), Inches(0.28),
                 size=10, bold=True, color=WHITE)
        anims.append((ht, 300))
        cx += cw2
    anims.append((thead, 300))

    test_rows = [
        ('RF-01', 'Inventario de almacenamiento', '12', 'OK'),
        ('RF-02', 'Configuración de políticas', '10', 'OK'),
        ('RF-03', 'Simulación', '11', 'OK'),
        ('RF-04', 'Limpieza FIFO producción', '14', 'OK'),
        ('RF-05', 'Bitácora y auditoría', '10', 'OK'),
        ('RF-07', 'Limpieza programada', '10', 'OK'),
        ('RF-08', 'Monitoreo preventivo', '10', 'OK'),
        ('Edge', 'Errores, permisos, concurrencia', '5', 'OK'),
        ('v2.0', 'RF-06 alarmas, RF-09 RBAC', '7', 'N/A'),
    ]
    for i, (req, desc, tests, estado) in enumerate(test_rows):
        ry = ty + Inches(0.42) + i * Inches(0.52)
        row_fill = GR100 if i % 2 == 0 else WHITE
        rc2 = rect(sl, tx, ry, tw, Inches(0.5), fill=row_fill, line_color=GR200)
        vals = [req, desc, tests, estado]
        cx = tx
        for v, cw2 in zip(vals, col_ws):
            bold_v = v == req
            col_v = B500 if req.startswith('RF') else GR500
            if v == estado:
                col_v = G500 if v == 'OK' else GR400
            vt = txb(sl, v, cx + Inches(0.1), ry + Inches(0.12), cw2 - Inches(0.1), Inches(0.3),
                     size=10, bold=bold_v, color=col_v if v == req else GR700)
            anims.append((vt, 350 + i*60))
            cx += cw2
        anims.append((rc2, 350 + i*60))

    add_logo(sl, dark_bg=False)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_10_instalacion(prs):
    """Slide 10: Plan de instalación — bg-white"""
    sl = new_slide(prs)
    set_bg(sl, WHITE)
    anims = []

    lb = add_label(sl, 'Implementación', color=B500)
    anims.append((lb, 0))
    tt = add_title(sl, 'Plan de instalación: 25 minutos', size=28, color=GR700)
    anims.append((tt, 100))
    sub = add_subtitle(sl, 'Sin reinicio del servidor. Sin modificaciones al sistema operativo.',
                       size=13, color=GR600)
    anims.append((sub, 200))

    steps = [
        ('1', 'Instalar .NET 8.0 Desktop Runtime', 'Descarga de Microsoft (~55 MB). Responsable: TI ODL', '5 min'),
        ('2', 'Crear carpeta C:\\FifoCleanup\\', 'Responsable: TI ODL', '1 min'),
        ('3', 'Copiar archivos publicados', '~50 MB. Responsable: IDC + TI ODL', '2 min'),
        ('4', 'Verificar permisos en carpeta de datos', 'Lectura/escritura/eliminación. Responsable: TI ODL', '2 min'),
        ('5', 'Ejecutar FifoCleanup + configurar', 'Establecer ruta de datos, umbral, cap. Responsable: IDC', '5 min'),
        ('6', 'Verificar inventario + simulación', 'Validar que detecta Assets correctamente. Responsable: IDC', '7 min'),
        ('7', 'Activar RF-07 y RF-08', 'Habilitar limpieza automática. Responsable: IDC + operador', '1 min'),
    ]
    lw = Inches(7.5)
    for i, (num, title_t, desc_t, time_t) in enumerate(steps):
        sy = MT + Inches(1.8) + i * Inches(0.72)
        sc = card(sl, ML, sy, lw, Inches(0.65), fill=GR100, border=GR200)
        num_c = rect(sl, ML + Inches(0.1), sy + Inches(0.1), Inches(0.42), Inches(0.42), fill=B500)
        num_t = txb(sl, num, ML + Inches(0.12), sy + Inches(0.12), Inches(0.38), Inches(0.35),
                    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        st = txb(sl, title_t, ML + Inches(0.65), sy + Inches(0.06), lw - Inches(1.6), Inches(0.28),
                 size=11, bold=True, color=GR700)
        sd = txb(sl, desc_t, ML + Inches(0.65), sy + Inches(0.36), lw - Inches(1.6), Inches(0.23),
                 size=9, color=GR500)
        tmt = txb(sl, time_t, ML + lw - Inches(1.1), sy + Inches(0.18), Inches(0.9), Inches(0.28),
                  size=10, bold=True, color=B500, align=PP_ALIGN.CENTER)
        d = 300 + i * 80
        anims.extend([(sc, d), (num_c, d), (num_t, d), (st, d), (sd, d), (tmt, d)])

    # Right: file structure
    rx = ML + lw + Inches(0.33)
    rw = Inches(3.75)
    rc2 = card(sl, rx, MT + Inches(1.8), rw, Inches(2.8), fill=GR100, border=GR300)
    rht = txb(sl, 'Lo que se instala', rx + Inches(0.15), MT + Inches(1.9),
              rw - Inches(0.3), Inches(0.28), size=11, bold=True, color=GR600)
    rft = txb_multi(sl, [
        'C:\\FifoCleanup\\',
        '  ├── FifoCleanup.exe',
        '  ├── *.dll (Engine, MVVM, Charts)',
        '  ├── *.runtimeconfig.json',
        '  └── fifo_config.json',
        '',
        'Total: < 50 MB en disco',
    ], rx + Inches(0.15), MT + Inches(2.25), rw - Inches(0.3), Inches(2.2),
        size=10, color=GR700, name='Consolas')
    anims.extend([(rc2, 400), (rht, 400), (rft, 420)])

    nc = card(sl, rx, MT + Inches(4.75), rw, Inches(2.0), fill=G100, border=G300)
    nht = txb(sl, '  Lo que NO se modifica', rx + Inches(0.15), MT + Inches(4.85),
              rw - Inches(0.3), Inches(0.28), size=11, bold=True, color=G500)
    no_list = ['Registro de Windows', 'Servicios del sistema', 'Puertos de red', 'Configuraciones del SO']
    for i, item in enumerate(no_list):
        it = txb(sl, f'--  {item}', rx + Inches(0.15), MT + Inches(5.22) + i * Inches(0.3),
                 rw - Inches(0.3), Inches(0.25), size=10, color=GR600)
        anims.append((it, 500 + i*60))
    anims.extend([(nc, 480), (nht, 490)])

    add_logo(sl, dark_bg=False)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_11_rollback(prs):
    """Slide 11: Rollback — bg-dark"""
    sl = new_slide(prs)
    set_bg(sl, B900)
    anims = []

    lb = add_label(sl, 'Plan de rollback', color=B400)
    anims.append((lb, 0))
    tt = add_title(sl, 'Desinstalación en 3 minutos', size=32, color=WHITE)
    anims.append((tt, 100))
    sub = add_subtitle(sl, 'Si algo sale mal, el servidor vuelve a su estado original sin rastro.',
                       size=14, color=B300)
    anims.append((sub, 200))

    steps = [
        ('close', 'Paso 1', 'Cerrar FifoCleanup.exe', '5 segundos'),
        ('trash', 'Paso 2', 'Eliminar carpeta C:\\FifoCleanup\\', '5 segundos'),
        ('cycle', 'Paso 3\n(opcional)', 'Desinstalar .NET 8.0 Runtime', '2 minutos'),
    ]
    sw2 = Inches(3.7)
    for i, (icon, step_t, desc_t, time_t) in enumerate(steps):
        sx = ML + i * (sw2 + Inches(0.16))
        sy = MT + Inches(2.0)
        sc = dark_card(sl, sx, sy, sw2, Inches(2.2))
        ico_sz = Inches(0.6)
        sit = add_icon(sl, icon, sx + (sw2 - ico_sz)/2, sy + Inches(0.22), ico_sz, ico_sz)
        sst = txb(sl, step_t, sx, sy + Inches(0.95), sw2, Inches(0.45),
                  size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        sdt = txb(sl, desc_t, sx + Inches(0.15), sy + Inches(1.4), sw2 - Inches(0.3), Inches(0.4),
                  size=11, color=GR400, align=PP_ALIGN.CENTER)
        stm = rect(sl, sx + Inches(0.7), sy + Inches(1.85), sw2 - Inches(1.4), Inches(0.28), fill=B800)
        stmt = txb(sl, time_t, sx + Inches(0.7), sy + Inches(1.87), sw2 - Inches(1.4), Inches(0.24),
                   size=10, bold=True, color=B400, align=PP_ALIGN.CENTER)
        d = 300 + i * 150
        anims.extend([(sc, d), (sit, d), (sst, d), (sdt, d), (stm, d), (stmt, d)])

    # 3 guarantee stats
    stats = [
        ('CERO', 'Impacto en monitoreo'),
        ('NINGUNO', 'Datos de monitoreo afectados'),
        ('NINGUNA', 'Configuración del servidor modificada'),
    ]
    gw = Inches(3.7)
    for i, (big_t, desc_t) in enumerate(stats):
        gx = ML + i * (gw + Inches(0.16))
        gy = MT + Inches(4.5)
        gc = dark_card(sl, gx, gy, gw, Inches(1.2))
        gc.line.color.rgb = G400
        gbt = txb(sl, big_t, gx, gy + Inches(0.15), gw, Inches(0.55),
                  size=26, bold=True, color=G400, align=PP_ALIGN.CENTER)
        gdt = txb(sl, desc_t, gx, gy + Inches(0.72), gw, Inches(0.35),
                  size=9, color=GR500, align=PP_ALIGN.CENTER)
        d = 600 + i * 100
        anims.extend([(gc, d), (gbt, d), (gdt, d)])

    add_logo(sl, dark_bg=True)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_12_documentacion(prs):
    """Slide 12: Documentación — bg-light"""
    sl = new_slide(prs)
    set_bg(sl, GR100)
    anims = []

    lb = add_label(sl, 'Entregables', color=B500)
    anims.append((lb, 0))
    tt = add_title(sl, '15+ documentos entregados en español', size=26, color=GR700)
    anims.append((tt, 100))

    docs = [
        ('document', 'README.md',             'Descripción, inicio rápido, estructura'),
        ('layers',   'Arquitectura',           'Diagrama de componentes y flujos'),
        ('gear',     'Especificación técnica', '17 parámetros, formatos, algoritmos'),
        ('log',      'Manual de operación',    'Operación diaria, interpretar Dashboard'),
        ('document', 'Runbook',                '14 procedimientos paso a paso'),
        ('check_circle','Troubleshooting',     '12 problemas comunes resueltos'),
        ('speech',   'FAQ',                    '30+ preguntas frecuentes'),
        ('download', 'Guía de instalación',    'Paso a paso con checklist'),
        ('gear',     'Guía de configuración',  'Detalle de cada parámetro'),
        ('play',     'Material de capacitación','Agenda 3h + ejercicios + evaluación'),
        ('document', 'Changelog + SLA',        'Versiones, severidades, tiempos de respuesta'),
        ('layers',   '7 ADRs + Requisitos',    'Decisiones, CU, RF, RNF, HU, CA'),
    ]

    cols = 4
    dw = Inches(2.8)
    dh = Inches(1.4)
    gap = Inches(0.13)
    dx0 = ML
    dy0 = MT + Inches(1.4)

    for i, (icon, dtitle, ddesc) in enumerate(docs):
        col = i % cols
        row = i // cols
        dx = dx0 + col * (dw + gap)
        dy = dy0 + row * (dh + gap)
        dc = card(sl, dx, dy, dw, dh, fill=WHITE, border=GR200)
        dit = add_icon(sl, icon, dx + Inches(0.15), dy + Inches(0.08), Inches(0.42), Inches(0.42))
        dtt = txb(sl, dtitle, dx + Inches(0.15), dy + Inches(0.55), dw - Inches(0.3), Inches(0.35),
                  size=11, bold=True, color=GR700)
        ddt = txb(sl, ddesc,  dx + Inches(0.15), dy + Inches(0.9),  dw - Inches(0.3), Inches(0.42),
                  size=9, color=GR500)
        d = 200 + i * 60
        anims.extend([(dc, d), (dit, d), (dtt, d), (ddt, d)])

    add_logo(sl, dark_bg=False)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_13_solicitud(prs):
    """Slide 13: Solicitud formal — bg-white"""
    sl = new_slide(prs)
    set_bg(sl, WHITE)
    anims = []

    lb = add_label(sl, 'Solicitud', color=B500)
    anims.append((lb, 0))
    tt = add_title(sl, 'Lo que necesitamos del equipo de TI', size=26, color=GR700)
    anims.append((tt, 100))

    cw = Inches(5.6)
    ly = MT + Inches(1.35)

    # Left: Solicitudes
    lx = ML
    lht = txb(sl, 'Solicitudes', lx, ly, cw, Inches(0.38),
              size=16, bold=True, color=B600)
    anims.append((lht, 200))

    requests = [
        ('1', 'Aprobación de instalación',
         'Instalar .NET 8.0 Runtime y FifoCleanup en SRVODLRTDNMICA'),
        ('2', 'Ventana de 30 minutos',
         'Para instalación y configuración inicial supervisada'),
        ('3', 'Acceso a carpeta de monitoreo',
         'Permisos de lectura, escritura y eliminación en la ruta de datos'),
        ('4', 'Contacto técnico',
         'Para coordinar la instalación y resolver dudas técnicas'),
    ]
    for i, (num, rtitle, rdesc) in enumerate(requests):
        ry = ly + Inches(0.5) + i * Inches(1.35)
        rc2 = card(sl, lx, ry, cw, Inches(1.22), fill=B100, border=B200)
        rnum = rect(sl, lx + Inches(0.15), ry + Inches(0.15), Inches(0.45), Inches(0.45), fill=B500)
        rnumt = txb(sl, num, lx + Inches(0.17), ry + Inches(0.17), Inches(0.4), Inches(0.38),
                    size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rtt = txb(sl, rtitle, lx + Inches(0.75), ry + Inches(0.1), cw - Inches(0.9), Inches(0.38),
                  size=13, bold=True, color=GR700)
        rdt = txb(sl, rdesc, lx + Inches(0.75), ry + Inches(0.52), cw - Inches(0.9), Inches(0.6),
                  size=10, color=GR600)
        d = 300 + i * 120
        anims.extend([(rc2, d), (rnum, d), (rnumt, d), (rtt, d), (rdt, d)])

    # Right: Compromisos
    rx = ML + cw + Inches(0.63)
    rw = CW - cw - Inches(0.63)
    rht = txb(sl, 'Nuestro compromiso', rx, ly, rw, Inches(0.38),
              size=16, bold=True, color=G500)
    anims.append((rht, 200))

    commitments = [
        'Instalación supervisada por ingeniero de IDC',
        'Capacitación 2–3 horas para operadores',
        'Soporte técnico post-go-live según SLA',
        'Documentación completa (15+ documentos)',
        'Suite de pruebas ejecutada (82/82 pasando)',
        'Plan de rollback inmediato (< 3 min)',
        'Roadmap v2.0 basado en feedback real',
    ]
    for i, comm in enumerate(commitments):
        cy = ly + Inches(0.5) + i * Inches(0.73)
        cc = card(sl, rx, cy, rw, Inches(0.63), fill=G100, border=G300)
        ct = txb(sl, f'--  {comm}', rx + Inches(0.15), cy + Inches(0.15), rw - Inches(0.3), Inches(0.35),
                 size=11, color=GR700)
        d = 300 + i * 80
        anims.extend([(cc, d), (ct, d)])

    add_logo(sl, dark_bg=False)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


def slide_14_resumen(prs):
    """Slide 14: Resumen ejecutivo — bg-gradient"""
    sl = new_slide(prs)
    set_bg(sl, B700)
    anims = []

    lb = add_label(sl, 'Resumen', color=B300)
    anims.append((lb, 0))
    tt = add_title(sl, 'FifoCleanup v1.0', size=36, color=WHITE)
    anims.append((tt, 100))
    sub = add_subtitle(sl, 'Gestión automática de espacio en disco para que el monitoreo nunca se detenga.',
                       size=15, color=B200)
    anims.append((sub, 200))

    # 4 stat cards
    stats4 = [('< 5%', 'CPU'), ('< 200 MB', 'RAM'), ('82/82', 'Tests pasando'), ('25 min', 'Instalación')]
    sw2 = Inches(2.8)
    sx0 = ML
    sy = MT + Inches(2.3)
    for i, (val, label_t) in enumerate(stats4):
        sc = dark_card(sl, sx0 + i*(sw2 + Inches(0.14)), sy, sw2, Inches(1.2))
        sv = txb(sl, val, sx0 + i*(sw2 + Inches(0.14)), sy + Inches(0.1), sw2, Inches(0.7),
                 size=26, bold=True, color=G400, align=PP_ALIGN.CENTER)
        sl2 = txb(sl, label_t, sx0 + i*(sw2 + Inches(0.14)), sy + Inches(0.82), sw2, Inches(0.3),
                  size=9, color=GR500, align=PP_ALIGN.CENTER)
        d = 300 + i * 100
        anims.extend([(sc, d), (sv, d), (sl2, d)])

    # 3 highlight cards
    highlights = [
        ('Bajo riesgo', 'Sin internet, sin puertos, sin admin, sin registro. Desinstalación en 3 min.'),
        ('Cero intervención', 'RF-07 + RF-08 trabajan 24/7 de forma autónoma. Auditoría completa.'),
        ('Documentado', '15+ documentos en español. Capacitación incluida. SLA definido.'),
    ]
    hw = Inches(3.7)
    hx0 = ML
    hy = MT + Inches(3.75)
    for i, (htitle, hdesc) in enumerate(highlights):
        hx = hx0 + i * (hw + Inches(0.16))
        hc = dark_card(sl, hx, hy, hw, Inches(1.65))
        htt = txb(sl, htitle, hx + Inches(0.2), hy + Inches(0.15), hw - Inches(0.4), Inches(0.38),
                  size=15, bold=True, color=WHITE)
        hdt = txb(sl, hdesc, hx + Inches(0.2), hy + Inches(0.6), hw - Inches(0.4), Inches(0.9),
                  size=10, color=B300)
        d = 500 + i * 120
        anims.extend([(hc, d), (htt, d), (hdt, d)])

    # Footer
    foot = txb_multi(sl, [
        {'text': 'IDC Ingeniería', 'size': 16, 'bold': True, 'color': WHITE},
        {'text': 'Marzo 2026', 'size': 10, 'color': GR500},
    ], ML, MT + Inches(5.65), CW, Inches(0.65), align=PP_ALIGN.CENTER)
    anims.append((foot, 700))

    add_logo(sl, dark_bg=True)
    add_transition_zoom(sl)
    add_animations(sl, anims)
    return sl


# ═══════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # Flujo: Portada → Índice → Contenido → Índice → Contenido → ... (hub-and-spoke)
    # Transiciones: portada=push, índice-inicial=fade, contenido=zoom-in, índice-retorno=zoom-out
    #
    # Total: 1 portada + 1 índice-inicial + 14×(contenido + índice-retorno) - 1 = 29 slides
    # (no hay índice después del último contenido)

    print("Generando diapositivas...")
    n = [1]
    def ok(label):
        print(f"  OK Slide {n[0]:2d}/29 - {label}")
        n[0] += 1

    slide_0_portada(prs);                       ok("Portada")
    slide_index(prs, first=True,  highlight=0); ok("Indice (inicial, proxima: El Problema)")

    slide_1_problema(prs);                      ok("El Problema")
    slide_index(prs, first=False, highlight=1); ok("Indice (retorno)")

    slide_2_solucion(prs);                      ok("La Solucion")
    slide_index(prs, first=False, highlight=2); ok("Indice (retorno)")

    slide_3_antes_despues(prs);                 ok("Antes vs. Despues")
    slide_index(prs, first=False, highlight=3); ok("Indice (retorno)")

    slide_4_arquitectura(prs);                  ok("Arquitectura")
    slide_index(prs, first=False, highlight=4); ok("Indice (retorno)")

    slide_5_fifo_dual(prs);                     ok("FIFO Dual")
    slide_index(prs, first=False, highlight=5); ok("Indice (retorno)")

    slide_6_dashboard(prs);                     ok("Dashboard")
    slide_index(prs, first=False, highlight=6); ok("Indice (retorno)")

    slide_7_impacto(prs);                       ok("Impacto")
    slide_index(prs, first=False, highlight=7); ok("Indice (retorno)")

    slide_8_seguridad(prs);                     ok("Seguridad")
    slide_index(prs, first=False, highlight=8); ok("Indice (retorno)")

    slide_9_pruebas(prs);                       ok("Pruebas")
    slide_index(prs, first=False, highlight=9); ok("Indice (retorno)")

    slide_10_instalacion(prs);                  ok("Instalacion")
    slide_index(prs, first=False, highlight=10);ok("Indice (retorno)")

    slide_11_rollback(prs);                     ok("Rollback")
    slide_index(prs, first=False, highlight=11);ok("Indice (retorno)")

    slide_12_documentacion(prs);                ok("Documentacion")
    slide_index(prs, first=False, highlight=12);ok("Indice (retorno)")

    slide_13_solicitud(prs);                    ok("Solicitud")
    slide_index(prs, first=False, highlight=13);ok("Indice (retorno)")

    slide_14_resumen(prs);                      ok("Resumen (fin)")

    out = r"c:\Users\IDC INGENIERIA\OneDrive\IDC\Proyectos\Ejecucion_Proyectos\03_Almacenamiento_FIFO\FifoCleanup_Presentacion_TI_v3.pptx"
    prs.save(out)
    print(f"\nListo! Guardado en:\n   {out}")

if __name__ == '__main__':
    main()
